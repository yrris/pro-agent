"""PPTX 生成脚本（python-pptx）：标题页 + 要点页序列 → presentation.pptx。"""

import json
import os
import sys


def main() -> None:
    args = json.loads(sys.argv[1]) if len(sys.argv) > 1 else {}
    title = str(args.get("title") or "演示文稿")
    subtitle = str(args.get("subtitle") or "")
    slides = list(args.get("slides") or [])
    if not slides:
        print("生成失败: slides 不能为空")
        sys.exit(1)

    from pptx import Presentation

    prs = Presentation()
    # 标题页
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = subtitle

    # 要点页（Title and Content 版式）
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(s.get("title") or "")
        body = slide.placeholders[1].text_frame
        bullets = [str(b) for b in (s.get("bullets") or [])]
        if bullets:
            body.text = bullets[0]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = b
                p.level = 0

    out_dir = os.environ.get("SKILL_OUTPUT_DIR", ".")
    path = os.path.join(out_dir, "presentation.pptx")
    prs.save(path)
    print(f"已生成 {len(slides) + 1} 页演示文稿 presentation.pptx（{os.path.getsize(path)} 字节）")


if __name__ == "__main__":
    main()
